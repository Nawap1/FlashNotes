import os
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from pptx import Presentation

def read_file(file_path):
    """
    Reads the content of a file and returns the text.
    Supported file types: .txt, .pdf, .pptx

    :param file_path: Path to the file to be read
    :return: Text content of the file
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"The file {file_path} does not exist.")

    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == '.txt':
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    elif file_extension == '.pdf':
        return read_pdf(file_path)

    elif file_extension == '.pptx':
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

    else:
        raise ValueError(f"Unsupported file extension: {file_extension}")

def read_pdf(file_path):
    """
    Reads the content of a PDF file and returns the text.
    If the PDF contains images, OCR is used to extract text.

    :param file_path: Path to the PDF file
    :return: Text content of the PDF
    """
    doc = fitz.open(file_path)
    text = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        extracted_text = page.get_text("text")
        if extracted_text.strip():
            text.append(extracted_text)
        else:
            # If no text is extracted, use OCR on the page
            text.append(extract_text_from_image_page(page))
    return '\n'.join(text)

def extract_text_from_image_page(page):
    """
    Extracts text from an image-based PDF page using OCR.

    :param page: PyMuPDF Page object
    :return: Text content of the page
    """
    zoom = 2.0  # Zoom factor to improve OCR accuracy
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat)
    image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

    # Use pytesseract to perform OCR
    text = pytesseract.image_to_string(image)
    return text
