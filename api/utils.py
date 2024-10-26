import os
import json
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from pptx import Presentation


class DocumentReader:
    """A class for reading and extracting text from various document formats."""

    def __init__(self, file_path):
        """
        Initialize the DocumentReader with a file path.

        Args:
            file_path (str): Path to the document file

        Raises:
            FileNotFoundError: If the specified file does not exist
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"The file {file_path} does not exist.")
            
        self.file_path = file_path
        self.file_extension = os.path.splitext(file_path)[1].lower()

    def extract_text(self):
        """
        Extract text content from the document based on its file type.

        Returns:
            str: Extracted text content from the document

        Raises:
            ValueError: If the file type is not supported
        """
        handlers = {
            '.txt': self._extract_text_from_txt,
            '.pdf': self._extract_text_from_pdf,
            '.pptx': self._extract_text_from_pptx
        }

        handler = handlers.get(self.file_extension)
        if handler:
            return handler()
        else:
            raise ValueError(f"Unsupported file extension: {self.file_extension}")

    def _extract_text_from_txt(self):
        """
        Extract text from a plain text file.

        Returns:
            str: Content of the text file
        """
        with open(self.file_path, 'r', encoding='utf-8') as file:
            return file.read()

    def _extract_text_from_pdf(self):
        """
        Extract text from a PDF file, including text from images using OCR.

        Returns:
            str: Combined text content from all PDF pages
        """
        doc = fitz.open(self.file_path)
        text_contents = []

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text("text")

            if text.strip():
                text_contents.append(text)
            else:
                # If no text is found, try OCR
                text_contents.append(self._extract_text_from_pdf_image(page))

        return '\n'.join(text_contents)

    def _extract_text_from_pdf_image(self, page):
        """
        Extract text from an image-based PDF page using OCR.

        Args:
            page: PyMuPDF Page object

        Returns:
            str: Extracted text from the page image
        """
        zoom = 2.0  # Higher zoom factor for better OCR accuracy
        matrix = fitz.Matrix(zoom, zoom)
        pixmap = page.get_pixmap(matrix=matrix)
        
        image = Image.frombytes(
            "RGB", 
            [pixmap.width, pixmap.height], 
            pixmap.samples
        )
        
        return pytesseract.image_to_string(image)

    def _extract_text_from_pptx(self):
        """
        Extract text from a PowerPoint presentation.

        Returns:
            str: Combined text content from all slides
        """
        presentation = Presentation(self.file_path)
        text_contents = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_contents.append(shape.text)

        return '\n'.join(text_contents)


def parse_json(json_string):
    """
    Parse a JSON string into a Python object.

    Args:
        json_string (str): JSON string to parse, optionally with markdown code block markers

    Returns:
        dict/list: Parsed JSON object

    Raises:
        json.JSONDecodeError: If the JSON string is invalid
    """
    # Remove markdown code block markers if present
    cleaned_string = json_string.replace('```json', '').replace('```', '').strip()
    
    return json.loads(cleaned_string)

